"""
Universal Document Reader - Nursing Training AI
Reads: PDF, DOCX, DOC, EPUB, TXT, RTF, PPTX, XLSX, HTML, CSV, MD, JSON, XML
"""

import os
import json
import csv
from pathlib import Path
from typing import Optional


class DocumentReader:
    """Reads text content from multiple document formats"""

    SUPPORTED_FORMATS = {
        '.pdf', '.docx', '.doc', '.epub', '.txt', '.rtf',
        '.pptx', '.xlsx', '.xls', '.html', '.htm',
        '.csv', '.md', '.json', '.xml', '.log', '.tsv'
    }

    @staticmethod
    def read(file_path: str) -> Optional[str]:
        """Read any supported document and return text content"""
        path = Path(file_path)
        ext = path.suffix.lower()

        readers = {
            '.pdf': DocumentReader._read_pdf,
            '.docx': DocumentReader._read_docx,
            '.doc': DocumentReader._read_doc,
            '.epub': DocumentReader._read_epub,
            '.txt': DocumentReader._read_text,
            '.md': DocumentReader._read_text,
            '.log': DocumentReader._read_text,
            '.csv': DocumentReader._read_csv,
            '.tsv': DocumentReader._read_csv,
            '.rtf': DocumentReader._read_rtf,
            '.pptx': DocumentReader._read_pptx,
            '.xlsx': DocumentReader._read_xlsx,
            '.xls': DocumentReader._read_xlsx,
            '.html': DocumentReader._read_html,
            '.htm': DocumentReader._read_html,
            '.json': DocumentReader._read_json,
            '.xml': DocumentReader._read_xml,
        }

        reader = readers.get(ext)
        if not reader:
            print(f"Unsupported format: {ext}")
            return None

        try:
            return reader(file_path)
        except Exception as e:
            print(f"Error reading {path.name}: {e}")
            return None

    @staticmethod
    def _read_pdf(path: str) -> str:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)

    @staticmethod
    def _read_docx(path: str) -> str:
        import docx
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def _read_doc(path: str) -> str:
        import mammoth
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
            return result.value

    @staticmethod
    def _read_epub(path: str) -> str:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(path)
        text = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text.append(soup.get_text())
        return "\n".join(text)

    @staticmethod
    def _read_text(path: str) -> str:
        import chardet
        with open(path, 'rb') as f:
            raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        with open(path, 'r', encoding=encoding, errors='replace') as f:
            return f.read()

    @staticmethod
    def _read_csv(path: str) -> str:
        import chardet
        with open(path, 'rb') as f:
            raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get('encoding', 'utf-8') or 'utf-8'

        lines = []
        delimiter = '\t' if path.endswith('.tsv') else ','
        with open(path, 'r', encoding=encoding, errors='replace') as f:
            reader = csv.reader(f, delimiter=delimiter)
            for row in reader:
                lines.append(" | ".join(row))
        return "\n".join(lines)

    @staticmethod
    def _read_rtf(path: str) -> str:
        from striprtf.striprtf import rtf_to_text
        with open(path, 'r', errors='replace') as f:
            return rtf_to_text(f.read())

    @staticmethod
    def _read_pptx(path: str) -> str:
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def _read_xlsx(path: str) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        text = []
        for ws in wb.worksheets:
            text.append(f"--- Sheet: {ws.title} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    text.append(row_text)
        return "\n".join(text)

    @staticmethod
    def _read_html(path: str) -> str:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            return soup.get_text()

    @staticmethod
    def _read_json(path: str) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)

    @staticmethod
    def _read_xml(path: str) -> str:
        from bs4 import BeautifulSoup
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            soup = BeautifulSoup(f.read(), 'xml')
            return soup.get_text()

    @staticmethod
    def scan_directory(directory: str, recursive: bool = True) -> dict:
        """Scan directory and return file counts by format"""
        stats = {"supported": 0, "unsupported": 0, "by_format": {}, "files": []}
        path = Path(directory)

        pattern = "**/*" if recursive else "*"
        for file in path.glob(pattern):
            if file.is_file():
                ext = file.suffix.lower()
                if ext in DocumentReader.SUPPORTED_FORMATS:
                    stats["supported"] += 1
                    stats["by_format"][ext] = stats["by_format"].get(ext, 0) + 1
                    stats["files"].append(str(file))
                elif ext:
                    stats["unsupported"] += 1

        return stats


# Quick test
if __name__ == "__main__":
    print("Supported formats:", sorted(DocumentReader.SUPPORTED_FORMATS))
    print(f"Total: {len(DocumentReader.SUPPORTED_FORMATS)} formats")
